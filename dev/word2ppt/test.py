# -*- coding: UTF-8 -*-
'''
@学习网站      ：https://www.python-office.com
@读者群     ：http://www.python4office.cn/wechat-group/
@作者  ：B站/抖音/微博/小红书/公众号，都叫：程序员晚枫，微信：CoderWanFeng
@代码日期    ：2024/5/18 18:48 
@本段代码的视频说明     ：
'''
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# 读取Word文档
doc = Document('test.docx')

# 创建新的PPT
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # 选择一个布局，6通常是标题和内容布局

# 遍历Word文档的每个段落
for para in doc.paragraphs:
    # 创建新幻灯片
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加内容到幻灯片（这里仅作示例，实际需要根据内容格式化）
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = para.text
    body_shape.text = ""  # 在PPT中可能不需要Word段落的全部内容

# 保存PPT
prs.save('output.pptx')